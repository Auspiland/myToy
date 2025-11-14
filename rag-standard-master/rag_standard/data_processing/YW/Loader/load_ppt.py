import os, re, sys
from pptx import Presentation

BASE_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), '../../../..')) # 상대경로로 root 설정 (Rag_Standard_v2)
sys.path.append(BASE_DIR)
from rag_standard.utils import setup_logger

global log
log = setup_logger.setup_logger('LoadPPT', 'load_ppt')

def clean_text(text):
    """ 마크다운 표에 적합하게 텍스트를 정리하는 함수 - 텍스트의 줄바꿈(\n)을 HTML 줄바꿈 태그(<br>)로 변환 """
    if not text:
        return ""
    return ' <br>'.join(text.split('\n'))

def sort_shapes_by_position(shapes):
    """
    PPT의 도형들을 위치(상단->하단, 좌->우)에 따라 정렬하는 함수
    시각적 순서대로 텍스트를 읽어오기 위한 목적
    
    Args:
    shapes: PPT 도형 객체 리스트
    
    Returns:
    list: 위치에 따라 정렬된 도형 리스트
    """
    try:
        # top, left 속성이 있고 None이 아닌 도형만 필터링
        valid_shapes = [
            shape for shape in shapes
            if hasattr(shape, 'top') and hasattr(shape, 'left') 
            and shape.top is not None and shape.left is not None
        ]
        # 도형을 top, left 값으로 정렬
        return sorted(valid_shapes, key=lambda s: (s.top or 0, s.left or 0))
    except Exception:
        # 정렬 실패시 원본 리스트 반환
        return shapes

class LoadPPT:
    """
    load_ppt만 호출하면 됨
    - extract_slide_content로 텍스트만 있는 버전과 테이블이 포함된 버전 추출
        <- extract_text_with_bullets로 글머리 기호와 들여쓰기가 포함된 텍스트 추출
        <- process_table로 테이블을 처리하여 일반 텍스트 또는 마크다운 형식으로 변환
    """
    def __init__(self):
        pass

    def load_ppt(self, file_path):
        """
        PPT 파일을 로드하고 내용을 추출하는 함수 >> Python-pptx 사용
        
        Args:
            file_path (str): PPT 파일 경로
            
        Returns:
            tuple: (파일명, 확장자, 수정일자, 생성일자, 텍스트만 추출한 딕셔너리, 텍스트와 테이블을 모두 추출한 딕셔너리, 전체 텍스트를 저장한 문자열)
            - 파일 로드 실패 시 False

        딕셔너리 형식:
            {
                '페이지 번호': 페이지 내용,
                '페이지 번호': 페이지 내용,
                ...
            }
        """
        file_name_without_ext, extension = os.path.basename(file_path).rsplit('.', 1) # 파일명과 확장자 분리

        try:
            prs = Presentation(file_path) # PPT 파일 로드

            try:
                # 파일의 메타데이터에서 수정일자와 생성일자 추출 시도, 메타데이터가 없으면 파일 시스템의 시간 정보 사용
                modified_date = prs.core_properties.modified.strftime("%Y-%m-%d %H:%M:%S") if hasattr(prs, 'core_properties') and prs.core_properties.modified else str(os.path.getmtime(file_path).strftime("%Y-%m-%d %H:%M:%S"))
                created_date = prs.core_properties.created.strftime("%Y-%m-%d %H:%M:%S") if hasattr(prs, 'core_properties') and prs.core_properties.created else str(os.path.getctime(file_path).strftime("%Y-%m-%d %H:%M:%S"))
                
            except (AttributeError, TypeError):
                # 메타데이터 접근 실패 시 파일 시스템의 시간 정보 사용
                modified_date = os.path.getmtime(file_path)
                created_date = os.path.getctime(file_path)

            # 슬라이드 처리를 위한 결과 리스트 초기화
            text_only_result = [] # 텍스트만 추출한 결과
            text_and_table_result = [] # 테이블 포함된 결과

            # 각 슬라이드 처리
            for slide in prs.slides:
                try:
                    # 슬라이드에서 텍스트와 테이블 내용 추출
                    text_only, text_and_table = self.extract_slide_content(slide, file_name_without_ext, extension)

                    # 추출된 내용이 있는 경우
                    if text_only or text_and_table:
                        text_only_result.append(text_only.strip())
                        text_and_table_result.append(text_and_table.strip())
                    else:
                        # 빈 슬라이드의 경우 빈 문자열 추가하여 동기화 유지
                        text_only_result.append("")
                        text_and_table_result.append("")
                        
                except Exception as e:
                    # 슬라이드 처리 실패 시 오 로그 기록
                    log.warning(f"Error processing slide: {e}")
                    # 에러 발생한 슬라이드도 빈 문자열 추가하여 동기화 유지
                    text_only_result.append("")
                    text_and_table_result.append("")

            # 모든 페이지의 텍스트 내용을 하나의 문자열로 합치기 (각 페이지 사이에 두 줄 개행)
            total_text_table = "\n\n".join(text_and_table_result)

            log.info(f"Completed loading ppt file: {os.path.basename(file_path)}")

            # 최종 결과 반환 (슬라이드 번호를 키로 하는 딕셔너리 형태)
            return (
                file_name_without_ext, extension.lower(), modified_date, created_date,
                {str(i+1): content for i, content in enumerate(text_only_result)},
                {str(i+1): content for i, content in enumerate(text_and_table_result)}, total_text_table
            )
                    
        except Exception as e:
            log.info(f"Cannot load PPT file: {file_path} due to {e}")
            return False
    
    def extract_slide_content(self, slide, file_name_without_ext, extension):
        """
        PPT 슬라이드에서 텍스트와 테이블 내용을 추출하는 함수
        텍스트만 있는 버전과 테이블이 포함된 버전을 모두 생성
        에러가 발생해도 두 버전이 동일한 내용을 포함하도록 보장
        
        Args:
            slide: PPT 슬라이드 객체
            file_name_without_ext: 확장자 제외한 파일명
            extension: 파일 확장자
        
        Returns:
            tuple: (텍스트만 추출한 내용, 테이블 포함된 내용)
        """
        content_pairs = []  # (텍스트 버전, 테이블 포함 버전) 쌍을 저장할 리스트

        try:
            # 슬라이드 내 도형들을 위치 순서대로 정렬
            shapes = sort_shapes_by_position(slide.shapes)
        except Exception as e:
            # 정렬 실패시 원본 순서 사용
            shapes = slide.shapes
            log.warning(f"Error sorting shapes: {e}")

        # 각 도형에 대한 처리
        for shape in shapes:
            try:
                # 그룹화된 도형 처리 (shape_type이 6인 경우)
                if getattr(shape, 'shape_type', None) == 6 and hasattr(shape, 'shapes'):
                    for item in shape.shapes:
                        shape_text = self.extract_text_with_bullets(item, file_name_without_ext, extension)
                        if shape_text:
                            content_pairs.append((shape_text, shape_text))

                # 일반 텍스트 프레임 처리
                elif hasattr(shape, 'text_frame'):
                    shape_text = self.extract_text_with_bullets(shape, file_name_without_ext, extension)
                    if shape_text:
                        content_pairs.append((shape_text, shape_text))

                # 테이블 처리
                elif hasattr(shape, 'table') and shape.table:
                    try:
                        # 텍스트 버전과 마크다운 버전의 테이블 생성
                        table_text = self.process_table(shape.table, False, file_name_without_ext, extension)
                        markdown_table = self.process_table(shape.table, True, file_name_without_ext, extension)
                        
                        if table_text:
                            content_pairs.append((table_text, markdown_table or table_text))
                            
                    except Exception as e:
                        log.warning(f"Error processing table in shape: {e}")
                        continue

            except Exception as e:
                # 테이블이 아닌 도형 처리 에러는 info 레벨로 기록
                if "shape does not contain a table" in str(e):
                    log.info(f"Non-critical shape processing info: {e}")
                else:
                    log.warning(f"Error processing shape: {e}")
                continue

        # 최종 결과 조합
        text_only = "\n".join(pair[0] for pair in content_pairs if pair[0]) # 텍스트만 있는 버전
        text_and_table = "\n".join(pair[1] for pair in content_pairs if pair[1]) # 테이블 포함 버전

        return text_only, text_and_table
    
    def extract_text_with_bullets(self, shape, file_name_without_ext, extension):
        """
        PPT 도형에서 글머리 기호와 들여쓰기가 포함된 텍스트를 추출하는 함수
        
        Args:
            shape: PPT 도형 객체
            file_name_without_ext: 확장자 제외한 파일명
            extension: 파일 확장자
        
        Returns:
            str: 추출된 텍스트 문자열
        """
        # 텍스트 프레임이 없는 도형은 빈 문자열 반환
        if not hasattr(shape, 'text_frame') or not shape.text_frame:
            return ""
        
        text_parts = []
        for paragraph in shape.text_frame.paragraphs:
            try:
                # 들여쓰기 레벨에 따른 탭 추가
                indent = '\t' * (paragraph.level if hasattr(paragraph, 'level') else 0)
                
                # 글머리 기호 처리
                bullet = '- '  # 기본 글머리 기
                try:
                    if hasattr(paragraph, 'bullet') and paragraph.bullet and hasattr(paragraph.bullet, 'character'):
                        bullet = paragraph.bullet.character + ' '
                except Exception:
                    pass
                
                # 문단 텍스트 추출 및 포맷팅
                paragraph_text = paragraph.text if paragraph.text else ""
                if paragraph_text.strip():
                    line_text = f"{indent}{bullet}{paragraph_text}"
                    text_parts.append(re.sub(r'\s{2,}', ' ', line_text))
                    
            except Exception as e:
                log.warning(f"Error processing paragraph: {e}")
                continue
                
        return '\n'.join(text_parts)
    
    def process_table(self, table, include_markdown, file_name_without_ext, extension):
        """
        PPT 테이블을 처리하여 일반 텍스트 또는 마크다운 형식으로 변환하는 함수
        
        Args:
            table: PPT 테이블 객체
            include_markdown: 마크다운 형식 포함 여부 (True/False)
            file_name_without_ext: 확장자 제외한 파일명
            extension: 파일 확장자
        
        Returns:
            str: 변환된 테이블 문자열
        """
        # 테이블에 행이 없으면 빈 문자열 반환
        if not hasattr(table, 'rows') or len(table.rows) == 0:
            return ""

        if not include_markdown:
            # 일반 텍스트 형식으로 처리
            rows_text = []
            for row in table.rows:
                row_cells = []
                for cell in row.cells:
                    try:
                        # 셀의 텍스트 추출
                        if hasattr(cell, 'text_frame') and cell.text_frame:
                            cell_text = cell.text_frame.text.strip()
                            row_cells.append(cell_text)
                    except Exception:
                        row_cells.append("")
                # 행의 셀들을 공백으로 구분하여 연결
                if row_cells:
                    rows_text.append(re.sub(r'\s{2,}', ' ', ' '.join(row_cells)))
            return '\n'.join(rows_text)

        # 마크다운 형식으로 처리
        try:
            # 헤더 행 처리
            header_cells = []
            for cell in table.rows[0].cells:
                try:
                    cell_text = clean_text(cell.text_frame.text) if hasattr(cell, 'text_frame') and cell.text_frame else ""
                    header_cells.append(cell_text)
                except Exception:
                    header_cells.append("")

            # 빈 테이블이면 빈 문자열 반환
            if not any(header_cells):
                return ""

            # 마크다운 테이블 형식 구성
            table_lines = [
                "| " + " | ".join(header_cells) + " |", # 헤더 행
                "|" + " --- |" * len(table.columns) # 구분선
            ]

            # 데이터 행 처리
            for row in table.rows[1:]:
                cell_contents = []
                for cell in row.cells:
                    try:
                        if hasattr(cell, 'text_frame') and cell.text_frame:
                            # 셀 텍스트 정리 (줄바꿈, 여러 공백 처리)
                            cell_text = clean_text(cell.text_frame.text)
                            cell_text = re.sub(r'<br> ', '', cell_text).rstrip(' <br>')
                            cell_text = re.sub(r'\s{2,}', ' ', cell_text)
                            cell_contents.append(cell_text)
                        else:
                            cell_contents.append("")
                    except Exception:
                        cell_contents.append("")

                # 비어있지 않은 행만 추가
                if any(cell_contents):
                    table_lines.append("| " + " | ".join(cell_contents) + " |")

            return '\n'.join(table_lines) + '\n'
            
        except Exception as e:
            # 마크다운 포맷팅 관련 에러는 info 레벨로 기록
            if "'<' not supported between instances of 'slice' and 'int'" in str(e):
                log.info(f"Non-critical table formatting info: {e}")
            else:
                log.warning(f"Error processing table: {e}")
            return ""